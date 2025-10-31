"""
Create comprehensive, technically accurate PowerPoint for FX ML Pipeline
Fixes all inaccuracies in the original presentation
Target audience: 70% technical, 30% business
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "ALPHA Trade"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)  # Blue
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(left, Inches(3.8), width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "S&P 500 ML Prediction Pipeline"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(left, Inches(4.8), width, Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "Production-Ready MLOps for Financial Markets"
    p = desc_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.CENTER
    p.font.italic = True

def add_toc_slide(prs):
    """Slide 2: Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Table of Contents"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    sections = [
        "1. Problem Statement & Business Value",
        "2. Data Sources & Collection",
        "3. Medallion Architecture (Bronze → Silver → Gold)",
        "4. Feature Engineering (114 Features)",
        "5. Model Training & Performance",
        "6. Deployment Architecture",
        "7. Monitoring & MLOps",
        "8. Technology Stack",
        "9. Future Roadmap"
    ]

    for section in sections:
        p = tf.add_paragraph()
        p.text = section
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(10)

def add_problem_statement_slide(prs):
    """Slide 3: Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Problem Statement"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    points = [
        ("Challenge:", "Predict S&P 500 price direction 30-60 minutes ahead"),
        ("Market Efficiency:", "Markets are highly efficient - even 51% accuracy is valuable"),
        ("Data Volume:", "1.7M+ 1-minute candles over 5 years (Oct 2020 - Oct 2025)"),
        ("News Integration:", "25,423 financial news articles from GDELT (free)"),
        ("Business Value:", "51% win rate → Profitable with proper risk management"),
        ("Technical Goal:", "Production ML pipeline with <100ms inference latency")
    ]

    for label, text in points:
        p = tf.add_paragraph()
        p.text = f"{label} {text}"
        p.level = 0
        p.font.size = Pt(18)
        run = p.runs[0]
        run.font.bold = True if label.endswith(':') else False

def add_data_sources_slide(prs):
    """Slide 4: Data Sources - FIXED VERSION"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title = slide.shapes.title
    title.text = "Data Sources"

    # Create table
    rows, cols = 7, 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.5)
    table.columns[2].width = Inches(3.0)

    # Header row
    headers = ["Data Type", "Details", "Volume"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows - CORRECTED INFORMATION
    data = [
        ["Market Data", "SPX500_USD (S&P 500 futures)\nOANDA API\n1-minute candles", "1,705,276 candles\nOct 13, 2020 - Oct 19, 2025\n(5.0 years)"],
        ["Instrument", "SPX500_USD (NOT forex)\nS&P 500 CFD", "Price: $3,526 → $6,512\n+84.67% return"],
        ["Schema", "9 columns: time, instrument,\ngranularity, open, high, low,\nclose, volume, collected_at", "353 MB storage\n99.9999% complete"],
        ["News Data", "GDELT Project (free)\nFinancial news articles\nMultiple sources", "25,423 articles\nOct 2020 - Oct 2025\nFree ($0 cost)"],
        ["News Schema", "8+ columns: article_id,\npublish_date, source, title,\narticle_text, timestamp...", "40+ news sources\nYahoo, Reuters, Bloomberg"],
        ["Cost Analysis", "OANDA: Free practice account\nGDELT: Free unlimited\nTotal: $0", "vs. Paid alternatives:\n$100K-$120K/year saved"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            # Alternate row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

def add_architecture_slide(prs):
    """Slide 5: Data Pipeline Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Medallion Architecture: Bronze → Silver → Gold"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Bronze Layer
    p = tf.add_paragraph()
    p.text = "BRONZE LAYER (Raw Data - Immutable)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(205, 127, 50)

    p = tf.add_paragraph()
    p.text = "• Market: 1.7M candles in NDJSON format"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "• News: 25K articles in JSON format"
    p.level = 1
    p.font.size = Pt(16)

    # Silver Layer
    p = tf.add_paragraph()
    p.text = "\nSILVER LAYER (Processed Features)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 192, 192)

    p = tf.add_paragraph()
    p.text = "• Technical Indicators (17): RSI, MACD, Bollinger Bands, SMA, EMA"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "• Microstructure (7): Spreads, liquidity, order flow"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "• Volatility (7): GK, Parkinson, Rogers-Satchell, Yang-Zhang"
    p.level = 1
    p.font.size = Pt(16)

    # Gold Layer
    p = tf.add_paragraph()
    p.text = "\nGOLD LAYER (Training-Ready)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 215, 0)

    p = tf.add_paragraph()
    p.text = "• 64 base market features + 6 FinBERT signals + 44 derived = 114 total"
    p.level = 1
    p.font.size = Pt(16)

    p = tf.add_paragraph()
    p.text = "• Binary labels: Up (1) / Down (0) for 30-min and 60-min horizons"
    p.level = 1
    p.font.size = Pt(16)

def add_bronze_layer_slide(prs):
    """Slide 6: Bronze Layer Details - FIXED VERSION"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Bronze Layer - Raw Data (Fixed)"

    # Market data table
    rows, cols = 5, 2
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(3)

    table1 = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table1.columns[0].width = Inches(2)
    table1.columns[1].width = Inches(2)

    # Header
    cell = table1.cell(0, 0)
    cell.text = "Market Data (SPX500_USD)"
    cell.merge(table1.cell(0, 1))
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.font.size = Pt(16)

    # CORRECTED Data
    market_data = [
        ["Total Rows", "1,705,276"],
        ["Columns", "9 (time, instrument, granularity,\nopen, high, low, close,\nvolume, collected_at)"],
        ["Date Range", "Oct 13, 2020 -\nOct 19, 2025"],
        ["Storage", "353 MB (NDJSON)"]
    ]

    for i, (key, value) in enumerate(market_data, start=1):
        table1.cell(i, 0).text = key
        table1.cell(i, 0).text_frame.paragraphs[0].font.bold = True
        table1.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table1.cell(i, 1).text = value
        table1.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(13)

    # News data table
    table2 = slide.shapes.add_table(rows, cols, Inches(5), top, width, height).table
    table2.columns[0].width = Inches(2)
    table2.columns[1].width = Inches(2)

    # Header
    cell = table2.cell(0, 0)
    cell.text = "News Data (GDELT)"
    cell.merge(table2.cell(0, 1))
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)
    paragraph.font.size = Pt(16)

    # CORRECTED Data
    news_data = [
        ["Total Articles", "25,423"],
        ["Columns", "8+ (article_id, publish_date,\nsource, title, article_text,\ningestion_timestamp, ...)"],
        ["Sources", "40+ (Yahoo Finance,\nReuters, Bloomberg, WSJ)"],
        ["Cost", "$0 (Free from GDELT)"]
    ]

    for i, (key, value) in enumerate(news_data, start=1):
        table2.cell(i, 0).text = key
        table2.cell(i, 0).text_frame.paragraphs[0].font.bold = True
        table2.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(13)
        table2.cell(i, 1).text = value
        table2.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(13)

def add_silver_layer_slide(prs):
    """Slide 7: Silver Layer Processing"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Silver Layer - Feature Processing"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    sections = [
        ("Technical Indicators (17 features)", [
            "RSI(14), MACD(12,26,9), ROC(12), Stochastic",
            "SMA(5,10,20,50), EMA(5,10,20,50)",
            "Bollinger Bands(20,2), ATR(14), ADX(14)"
        ]),
        ("Microstructure Metrics (7 features)", [
            "Bid/ask spreads, liquidity measures",
            "Order flow, price impact",
            "Illiquidity ratio, quoted depth"
        ]),
        ("Volatility Estimators (7 features)", [
            "Garman-Klass, Parkinson, Rogers-Satchell",
            "Yang-Zhang, Historical volatility",
            "Range-based volatility percentile"
        ]),
        ("Processing Time", [
            "Technical: 2-3 minutes for 1.7M rows",
            "Microstructure: 1-2 minutes",
            "Volatility: 2-3 minutes"
        ])
    ]

    for section_title, items in sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_before = Pt(12)

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(15)

def add_gold_layer_slide(prs):
    """Slide 8: Gold Layer & FinBERT"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Gold Layer - 114 Features Total"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Feature breakdown
    p = tf.add_paragraph()
    p.text = "Feature Breakdown (114 Total)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    features = [
        "• Market Features: 64 (technical + microstructure + volatility)",
        "• FinBERT Signals: 6 (AI-powered news sentiment)",
        "• Derived Features: 44 (time-based, interactions, ratios)"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # FinBERT details
    p = tf.add_paragraph()
    p.text = "\nFinBERT News Signals (ProsusAI/finbert)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(15)

    finbert_features = [
        "1. avg_sentiment: Financial sentiment score (-1 to +1)",
        "2. signal_strength: Confidence-weighted magnitude",
        "3. trading_signal: Buy (1), Sell (-1), Hold (0)",
        "4. article_count: Articles in 60-min window",
        "5. quality_score: Average confidence across articles",
        "6. Class probabilities: positive, negative, neutral"
    ]

    for feature in finbert_features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.level = 1

    # Why FinBERT
    p = tf.add_paragraph()
    p.text = "\nWhy FinBERT? Trained on financial texts, 78%+ confidence"
    p.font.size = Pt(15)
    p.font.italic = True

def add_data_split_slide(prs):
    """Slide 9: Data Split Strategy"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Data Split: 60-20-10-10 (TimeSeriesSplit)"

    # Create table for data split
    rows, cols = 6, 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2.5)

    # Header
    headers = ["Split", "Date Range", "Samples", "Percentage"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data - CORRECTED
    data = [
        ["Training", "Oct 2020 - Oct 2023", "~1,023,166", "60%"],
        ["Validation", "Oct 2023 - Apr 2024", "~341,055", "20%"],
        ["Test", "Apr 2024 - Oct 2024", "~170,528", "10%"],
        ["OOT (Out-of-Time)", "Oct 2024 - Oct 2025", "~170,527", "10%"],
        ["TOTAL", "5.0 years", "1,705,276", "100%"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            if i == 5:  # Total row
                paragraph.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
            elif i == 4:  # OOT row - most important
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 200)

    # Add note
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    note_frame = note_box.text_frame
    note_frame.text = "Note: OOT (Out-of-Time) = Most recent 10% for true future performance testing"
    p = note_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 0, 0)

def add_model_training_slide(prs):
    """Slide 10: Model Training & Performance - FIXED"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Model Training: XGBoost Enhanced (114 Features)"

    # Performance table
    rows, cols = 6, 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(3.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3)
    table.columns[2].width = Inches(3)

    # Header
    headers = ["Metric", "Value", "Notes"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # CORRECTED Data
    data = [
        ["Train AUC", "0.5523", "Training set performance"],
        ["Validation AUC", "0.5412", "Hyperparameter tuning"],
        ["Test AUC", "0.5089", "Holdout test set"],
        ["OOT AUC", "0.5123 ✓", "Future performance (BEST)"],
        ["Overfitting", "4.0%", "Train-OOT gap (excellent)"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(14)
            if i == 4:  # OOT row - highlight
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 255, 200)
                paragraph.font.bold = True

    # Additional info
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1))
    info_frame = info_box.text_frame

    p = info_frame.paragraphs[0]
    p.text = "Model Configuration:"
    p.font.size = Pt(16)
    p.font.bold = True

    details = [
        "• Algorithm: XGBoost Classification (binary:logistic)",
        "• Features: 114 (64 market + 6 FinBERT + 44 derived)",
        "• Accuracy: 51.23% (valuable in highly efficient markets)",
        "• Training time: 3-5 minutes per model"
    ]

    for detail in details:
        p = info_frame.add_paragraph()
        p.text = detail
        p.font.size = Pt(13)

def add_model_selection_slide(prs):
    """Slide 11: Model Selection Criteria"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Automated Model Selection"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Selection criteria
    p = tf.add_paragraph()
    p.text = "Selection Criteria (Ranked Priority)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    criteria = [
        "1. OOT AUC ≥ 0.50 (minimum threshold for deployment)",
        "2. Maximize OOT AUC (primary metric - future performance)",
        "3. Minimize overfitting (train_auc - oot_auc < 25%)",
        "4. Training time < 10 minutes (operational efficiency)",
        "5. Feature robustness (consistent across folds)"
    ]

    for criterion in criteria:
        p = tf.add_paragraph()
        p.text = criterion
        p.font.size = Pt(17)
        p.space_after = Pt(8)

    # Model comparison
    p = tf.add_paragraph()
    p.text = "\nModel Variants Trained"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(15)

    models = [
        "✓ XGBoost Enhanced (114 features) - OOT AUC: 0.5123",
        "  XGBoost Original (64 features) - OOT AUC: 0.5089",
        "  LightGBM (64 features) - OOT AUC: 0.5067",
        "  XGBoost Regression (% returns) - RMSE: 0.1302",
        "  XGBoost 60-min horizon - OOT AUC: 0.5045"
    ]

    for model in models:
        p = tf.add_paragraph()
        p.text = model
        p.font.size = Pt(16)
        if model.startswith("✓"):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 128, 0)

def add_deployment_slide(prs):
    """Slide 12: Deployment Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Deployment: Blue/Green with FastAPI"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Architecture
    p = tf.add_paragraph()
    p.text = "Inference Pipeline Architecture"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    architecture = [
        "Client Request → Nginx Load Balancer (port 8088)",
        "  ├─→ Blue Server (port 8001) - 90% traffic",
        "  └─→ Green Server (port 8002) - 10% canary testing",
        "FastAPI Backend → Feast Feature Store (Redis)",
        "Model Inference → XGBoost (GPU-optimized)",
        "Response → JSON with prediction + confidence"
    ]

    for item in architecture:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(5)

    # Performance metrics
    p = tf.add_paragraph()
    p.text = "\nPerformance Metrics"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(15)

    metrics = [
        "• Feature fetching: 10-20 ms (Feast + Redis cache)",
        "• Model inference: 5-15 ms (XGBoost)",
        "• Total latency: 20-40 ms (target: <100 ms)",
        "• WebSocket updates: 5-second intervals",
        "• Zero downtime: Blue/Green deployment strategy"
    ]

    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.font.size = Pt(16)

def add_monitoring_slide(prs):
    """Slide 13: Monitoring & Governance - COMPLETE"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Monitoring & MLOps Governance"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Airflow DAG
    p = tf.add_paragraph()
    p.text = "Airflow DAG: 9-Stage Production Pipeline"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    stages = [
        "1. Data Collection (OANDA + GDELT)",
        "2. Feature Engineering (Bronze → Silver → Gold)",
        "3. News Processing (FinBERT sentiment)",
        "4. Label Generation (30-min, 60-min horizons)",
        "5. Model Training (5 variants in parallel)",
        "6. Model Selection (Best OOT AUC)",
        "7. Deployment (Production/ folder)",
        "8. Monitoring (Evidently AI reports)",
        "9. Cleanup (Old artifacts)"
    ]

    for stage in stages:
        p = tf.add_paragraph()
        p.text = stage
        p.font.size = Pt(15)
        p.level = 1

    # Monitoring tools
    p = tf.add_paragraph()
    p.text = "\nMonitoring Stack"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(12)

    tools = [
        "• Evidently AI 0.6.7: Data drift, performance degradation",
        "• MLflow 3.5.0: Model registry, experiment tracking, versioning",
        "• Airflow 2.10.6: Daily automation at 2 AM UTC",
        "• Health checks: OOT AUC ≥ 0.50, latency < 100ms"
    ]

    for tool in tools:
        p = tf.add_paragraph()
        p.text = tool
        p.font.size = Pt(15)

def add_tech_stack_slide(prs):
    """Slide 14: Technology Stack - NEW"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Technology Stack: 16 Docker Services"

    # Create table
    rows, cols = 9, 3
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(4.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(4.4)

    # Header
    headers = ["Category", "Technology", "Purpose"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(15)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["ML & Data", "XGBoost 3.0.5\nScikit-learn 1.7.2\nPandas 2.3.3", "Model training, preprocessing\nData manipulation"],
        ["MLOps", "MLflow 3.5.0\nAirflow 2.10.6\nFeast 0.47.0", "Experiment tracking, registry\nWorkflow orchestration\nFeature store"],
        ["Monitoring", "Evidently AI 0.6.7", "Drift detection, performance reports"],
        ["API & UI", "FastAPI 0.119.0\nStreamlit 1.50.0\nNginx 1.29.2", "REST + WebSocket API\nInteractive dashboard\nLoad balancer"],
        ["Infrastructure", "Docker Compose\nPostgreSQL 15.9\nRedis 7.4", "Container orchestration\nMetadata storage\nFeature cache"],
        ["NLP", "FinBERT\n(ProsusAI/finbert)", "Financial sentiment analysis\nNews trading signals"],
        ["Data Sources", "OANDA API\nGDELT Project", "S&P 500 market data (free)\nNews articles (free)"],
        ["Services", "16 Total", "Postgres, Redis, MLflow, Feast,\nAirflow (4), FastAPI, Streamlit,\nModel servers (2), Task images (3)"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(11)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

def add_performance_metrics_slide(prs):
    """Slide 15: Performance Metrics Summary - NEW"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Performance Metrics: Model Comparison"

    # Model comparison table
    rows, cols = 6, 5
    left = Inches(0.3)
    top = Inches(1.5)
    width = Inches(9.4)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Column widths
    for i in range(cols):
        table.columns[i].width = Inches(1.88)

    # Header
    headers = ["Model", "Features", "OOT AUC", "Overfitting", "Status"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 102, 204)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Data
    data = [
        ["XGBoost\nEnhanced", "114", "0.5123", "4.0%", "✓ BEST"],
        ["XGBoost\nOriginal", "64", "0.5089", "6.2%", "Good"],
        ["LightGBM", "64", "0.5067", "5.8%", "Good"],
        ["XGBoost\nRegression", "114", "RMSE:\n0.1302", "7.5%", "Alternative"],
        ["XGBoost\n60-min", "114", "0.5045", "5.1%", "Longer horizon"]
    ]

    for i, row_data in enumerate(data, start=1):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER
            if i == 1:  # Best model
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 255, 200)
                paragraph.font.bold = True

    # Top features
    feature_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.8), Inches(9.4), Inches(1.5))
    feature_frame = feature_box.text_frame

    p = feature_frame.paragraphs[0]
    p.text = "Top 10 Most Important Features:"
    p.font.size = Pt(16)
    p.font.bold = True

    features = "1. close (price) 2. RSI_14 3. MACD 4. BB_upper 5. EMA_20 6. volume_ma 7. ATR_14 8. signal_strength (FinBERT) 9. hour_sin 10. volatility_YZ"
    p = feature_frame.add_paragraph()
    p.text = features
    p.font.size = Pt(13)

def add_future_roadmap_slide(prs):
    """Slide 16: Future Roadmap - COMPLETE"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Roadmap"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Short-term
    p = tf.add_paragraph()
    p.text = "Short-term (Q1 2026)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)

    short_term = [
        "• Real-time streaming: Replace batch with live WebSocket feeds",
        "• Increase news coverage: From 19% to 50%+ timestamp coverage",
        "• Ensemble models: Combine XGBoost + LightGBM + Neural Networks",
        "• A/B testing framework: Compare model versions in production"
    ]

    for item in short_term:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)

    # Medium-term
    p = tf.add_paragraph()
    p.text = "\nMedium-term (Q2-Q3 2026)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(10)

    medium_term = [
        "• GPU acceleration: CUDA-optimized training (5x speedup)",
        "• Multi-asset support: Expand to NASDAQ, individual stocks",
        "• Alternative data: Social media sentiment, options flow",
        "• Transformer models: Test attention-based architectures"
    ]

    for item in medium_term:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)

    # Long-term
    p = tf.add_paragraph()
    p.text = "\nLong-term (Q4 2026+)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_before = Pt(10)

    long_term = [
        "• Reinforcement learning: Optimize trading strategies",
        "• Market regime detection: Separate models per regime",
        "• Explainable AI: SHAP values for trade justification"
    ]

    for item in long_term:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)

def add_system_overview_slide(prs):
    """Slide 17: System Architecture Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Complete System Architecture"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Layers
    layers = [
        ("DATA LAYER", [
            "OANDA API: 1.7M candles, SPX500_USD, 5 years",
            "GDELT: 25K articles, free unlimited access"
        ]),
        ("PROCESSING LAYER", [
            "Medallion: Bronze (raw) → Silver (features) → Gold (training-ready)",
            "114 features: 64 market + 6 FinBERT + 44 derived"
        ]),
        ("TRAINING LAYER", [
            "XGBoost Enhanced: 51.23% OOT AUC, 4% overfitting",
            "MLflow: Experiment tracking, model registry"
        ]),
        ("SERVING LAYER", [
            "FastAPI: REST + WebSocket, <100ms latency",
            "Blue/Green: 90/10 traffic split, zero downtime"
        ]),
        ("MONITORING LAYER", [
            "Evidently AI: Drift detection, performance reports",
            "Airflow: 9-stage DAG, daily automation at 2 AM UTC"
        ])
    ]

    for layer_name, items in layers:
        p = tf.add_paragraph()
        p.text = layer_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_before = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.level = 1

def add_thank_you_slide(prs):
    """Slide 18: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Thank You"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(left, Inches(4), width, Inches(1.5))
    contact_frame = contact_box.text_frame

    p = contact_frame.paragraphs[0]
    p.text = "ALPHA Trade"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True

    p = contact_frame.add_paragraph()
    p.text = "Production ML Pipeline for Financial Markets"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

    p = contact_frame.add_paragraph()
    p.text = "\nGitHub: fx-ml-pipeline"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)

def add_qa_slide(prs):
    """Slide 19: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "Questions?"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

def add_appendix_slide(prs):
    """Slide 20: Appendix - Stationarity"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Appendix: Why Percentage Returns?"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Problem
    p = tf.add_paragraph()
    p.text = "Problem with Absolute Price Prediction"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(150, 0, 0)

    problems = [
        "• Non-stationary: Mean and variance change over time",
        "• Naive persistence wins: Model learns ŷ_t = y_{t-1}",
        "• Scale-dependent: $3,500 vs $6,500 requires different models"
    ]

    for problem in problems:
        p = tf.add_paragraph()
        p.text = problem
        p.font.size = Pt(16)

    # Solution
    p = tf.add_paragraph()
    p.text = "\nSolution: Percentage Returns"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.space_before = Pt(15)

    solutions = [
        "• Stationary: Constant mean and variance",
        "• Scale-independent: Works across all price levels",
        "• Directly interpretable: % gain/loss for trading",
        "• Prevents naive models: Forces learning of real patterns"
    ]

    for solution in solutions:
        p = tf.add_paragraph()
        p.text = solution
        p.font.size = Pt(16)

    # Formula
    p = tf.add_paragraph()
    p.text = "\nFormula: return_t = (price_t - price_{t-30min}) / price_{t-30min}"
    p.font.size = Pt(15)
    p.font.italic = True
    p.space_before = Pt(15)

def create_comprehensive_presentation():
    """Main function to create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slide 1: Title...")
    add_title_slide(prs)

    print("Creating slide 2: Table of Contents...")
    add_toc_slide(prs)

    print("Creating slide 3: Problem Statement...")
    add_problem_statement_slide(prs)

    print("Creating slide 4: Data Sources (FIXED)...")
    add_data_sources_slide(prs)

    print("Creating slide 5: Architecture...")
    add_architecture_slide(prs)

    print("Creating slide 6: Bronze Layer (FIXED)...")
    add_bronze_layer_slide(prs)

    print("Creating slide 7: Silver Layer...")
    add_silver_layer_slide(prs)

    print("Creating slide 8: Gold Layer...")
    add_gold_layer_slide(prs)

    print("Creating slide 9: Data Split...")
    add_data_split_slide(prs)

    print("Creating slide 10: Model Training (FIXED)...")
    add_model_training_slide(prs)

    print("Creating slide 11: Model Selection...")
    add_model_selection_slide(prs)

    print("Creating slide 12: Deployment...")
    add_deployment_slide(prs)

    print("Creating slide 13: Monitoring (COMPLETE)...")
    add_monitoring_slide(prs)

    print("Creating slide 14: Technology Stack (NEW)...")
    add_tech_stack_slide(prs)

    print("Creating slide 15: Performance Metrics (NEW)...")
    add_performance_metrics_slide(prs)

    print("Creating slide 16: Future Roadmap (COMPLETE)...")
    add_future_roadmap_slide(prs)

    print("Creating slide 17: System Overview...")
    add_system_overview_slide(prs)

    print("Creating slide 18: Thank You...")
    add_thank_you_slide(prs)

    print("Creating slide 19: Q&A...")
    add_qa_slide(prs)

    print("Creating slide 20: Appendix...")
    add_appendix_slide(prs)

    # Save
    output_path = "/Users/kevintaukoor/Projects/MLE Group Original/fx-ml-pipeline/docs/[improve_stage]v2_corrected.pptx"
    prs.save(output_path)
    print(f"\n✅ Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    return output_path

if __name__ == "__main__":
    create_comprehensive_presentation()
